from fastapi import FastAPI, HTTPException, UploadFile, File as FastAPIFile
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, Response
from pydantic import BaseModel
from databricks.sdk import WorkspaceClient
from databricks.sdk.service.files import FileInfo
from databricks.sdk.service.sql import StatementState
import os
import yaml
from dotenv import load_dotenv
from typing import List, Dict, Any, Optional
import tempfile
import shutil
import json
import fitz
import io
import atexit
import re
from abc import ABC, abstractmethod
import mimetypes

# New imports for format support
try:
    import openpyxl
    import xlsxwriter
    from openpyxl import Workbook
    from openpyxl.utils.dataframe import dataframe_to_rows
    EXCEL_SUPPORT = True
except ImportError:
    EXCEL_SUPPORT = False
    print("Excel support not available - install openpyxl and xlsxwriter")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    POWERPOINT_SUPPORT = True
except ImportError:
    POWERPOINT_SUPPORT = False
    print("PowerPoint support not available - install python-pptx")

try:
    import markdown
    import mistune
    MARKDOWN_SUPPORT = True
except ImportError:
    MARKDOWN_SUPPORT = False
    print("Markdown support not available - install markdown and mistune")

def load_yaml_config():
    """Load configuration from app.yaml file"""
    try:
        with open('app.yaml', 'r') as file:
            config = yaml.safe_load(file)
            # Convert env array to a dictionary for easy access
            yaml_config = {}
            if 'env' in config:
                for env_var in config['env']:
                    yaml_config[env_var['name']] = env_var['value']
            return yaml_config
    except Exception as e:
        print(f"Warning: Could not load app.yaml config: {e}")
        return {}

# Load YAML configuration
YAML_CONFIG = load_yaml_config()

# Utility functions for format detection and handling
def get_file_extension(filename: str) -> str:
    """Get file extension from filename"""
    return os.path.splitext(filename.lower())[1]

def detect_file_format(file_path: str) -> str:
    """Detect file format from path and return format type"""
    ext = get_file_extension(file_path)
    
    format_mapping = {
        '.md': 'markdown',
        '.markdown': 'markdown',
        '.xlsx': 'excel',
        '.xls': 'excel',
        '.pptx': 'powerpoint',
        '.ppt': 'powerpoint',
        '.pdf': 'pdf',
        '.txt': 'text',
        '.doc': 'word',
        '.docx': 'word',
        '.csv': 'csv',
        '.json': 'json'
    }
    
    return format_mapping.get(ext, 'unknown')

def validate_file_format(filename: str, supported_formats: list) -> bool:
    """Validate if file format is supported"""
    ext = get_file_extension(filename)
    return ext in supported_formats

def generate_output_filename(original_path: str, suffix: str, new_extension: str = None) -> str:
    """Generate output filename with suffix"""
    base_path = os.path.dirname(original_path)
    filename = os.path.splitext(os.path.basename(original_path))[0]
    ext = new_extension or os.path.splitext(original_path)[1]
    return os.path.join(base_path, f"{filename}_{suffix}{ext}").replace('\\', '/')

def get_mime_type(file_extension: str) -> str:
    """Get MIME type for file extension"""
    mime_mapping = {
        '.pdf': 'application/pdf',
        '.md': 'text/markdown',
        '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        '.xls': 'application/vnd.ms-excel',
        '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        '.ppt': 'application/vnd.ms-powerpoint',
        '.txt': 'text/plain',
        '.json': 'application/json',
        '.csv': 'text/csv'
    }
    return mime_mapping.get(file_extension.lower(), 'application/octet-stream')

# Supported formats configuration
SUPPORTED_UPLOAD_FORMATS = ['.pdf', '.md', '.xlsx', '.xls', '.pptx', '.ppt', '.txt', '.doc', '.docx', '.csv', '.json']
SUPPORTED_EXPORT_FORMATS = ['md', 'xlsx', 'pptx', 'pdf']

load_dotenv()

app = FastAPI()

# Configure CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Pydantic models
class ParseRequest(BaseModel):
    text: str


# Add new Pydantic models after existing ones
class WriteToTableRequest(BaseModel):
    file_paths: List[str]
    limit: int = 10

class QueryDeltaTableRequest(BaseModel):
    file_paths: List[str] = []
    limit: int = 10

# Helper functions
def get_uc_volume_path() -> str:
    """Get the current UC Volume path"""
    return current_volume_path or "/Volumes/main/default/ai_functions_demo"

def get_delta_table_path() -> str:
    """Get the current Delta table path"""  
    return current_delta_table_path or "main.default.ai_functions_demo_documents"

# Initialize Databricks client - uses automatic authentication in Databricks Apps
try:
    w = WorkspaceClient()  # Automatic authentication
    warehouse_id = os.getenv("DATABRICKS_WAREHOUSE_ID", YAML_CONFIG.get("DATABRICKS_WAREHOUSE_ID"))
    print(f"✅ Databricks client initialized with warehouse: {warehouse_id}")
except Exception as e:
    print(f"⚠️ Databricks client initialization failed: {e}")
    w = None
    warehouse_id = None

# Global variables to store dynamic configuration
current_warehouse_id = warehouse_id
current_volume_path = os.getenv("DATABRICKS_VOLUME_PATH", YAML_CONFIG.get("DATABRICKS_VOLUME_PATH"))
current_delta_table_path = os.getenv("DATABRICKS_DELTA_TABLE_PATH", YAML_CONFIG.get("DATABRICKS_DELTA_TABLE_PATH"))

# Base format handler class
class BaseFormatHandler(ABC):
    @abstractmethod
    def can_handle(self, file_path: str, mime_type: str = None) -> bool:
        """Check if handler can process this file type"""
        pass
    
    @abstractmethod
    def extract_content(self, file_path: str) -> str:
        """Extract text content from file"""
        pass
    
    @abstractmethod
    def redact_content(self, file_path: str, replacements: Dict[str, str]) -> str:
        """Redact content and return new file path"""
        pass
    
    @abstractmethod
    def get_supported_extensions(self) -> List[str]:
        """Return list of supported file extensions"""
        pass

# Markdown handler implementation
class MarkdownHandler(BaseFormatHandler):
    def can_handle(self, file_path: str, mime_type: str = None) -> bool:
        return (file_path.lower().endswith(('.md', '.markdown')) or 
                mime_type == 'text/markdown') and MARKDOWN_SUPPORT
    
    def extract_content(self, file_path: str) -> str:
        """Extract text content from Markdown file"""
        if not w:
            raise Exception("Databricks connection not configured")
        
        try:
            # Download file from UC Volume
            file_response = w.files.download(file_path=file_path)
            
            # Read content
            if hasattr(file_response, 'content'):
                content = file_response.content.decode('utf-8')
            elif hasattr(file_response, 'read'):
                content = file_response.read().decode('utf-8')
            else:
                content = str(file_response)
            
            return content
        except Exception as e:
            print(f"Error extracting Markdown content: {e}")
            return ""
    
    def redact_content(self, file_path: str, replacements: Dict[str, str]) -> str:
        """Redact Markdown content using text replacement"""
        if not replacements:
            return file_path
        
        try:
            # Download original file
            content = self.extract_content(file_path)
            
            # Perform text replacements
            redacted_content = content
            for search_text, replace_text in replacements.items():
                redacted_content = re.sub(re.escape(search_text), replace_text, redacted_content, flags=re.IGNORECASE)
            
            # Generate redacted filename
            redacted_path = generate_output_filename(file_path, "redacted", ".md")
            
            # Upload redacted content back to UC Volume
            with tempfile.NamedTemporaryFile(mode='w', encoding='utf-8', delete=False, suffix='.md') as temp_file:
                temp_file.write(redacted_content)
                temp_file_path = temp_file.name
            
            try:
                with open(temp_file_path, 'rb') as f:
                    w.files.upload(
                        file_path=redacted_path,
                        contents=f,
                        overwrite=True
                    )
            finally:
                os.unlink(temp_file_path)
            
            return redacted_path
        except Exception as e:
            print(f"Error redacting Markdown content: {e}")
            raise e
    
    def get_supported_extensions(self) -> List[str]:
        return ['.md', '.markdown']

# Excel handler implementation
class ExcelHandler(BaseFormatHandler):
    def can_handle(self, file_path: str, mime_type: str = None) -> bool:
        return (file_path.lower().endswith(('.xlsx', '.xls')) or 
                mime_type in ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet']) and EXCEL_SUPPORT
    
    def extract_content(self, file_path: str) -> str:
        """Extract text content from Excel sheets"""
        if not w:
            raise Exception("Databricks connection not configured")
        
        try:
            # Download file from UC Volume
            file_response = w.files.download(file_path=file_path)
            
            # Create temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as temp_file:
                if hasattr(file_response, 'content'):
                    temp_file.write(file_response.content)
                elif hasattr(file_response, 'read'):
                    temp_file.write(file_response.read())
                else:
                    temp_file.write(file_response)
                temp_file_path = temp_file.name
            
            try:
                # Read Excel file
                workbook = openpyxl.load_workbook(temp_file_path, data_only=True)
                content_parts = []
                
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    content_parts.append(f"## Sheet: {sheet_name}\n")
                    
                    for row in sheet.iter_rows(values_only=True):
                        row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                        if row_text.strip():
                            content_parts.append(row_text)
                    content_parts.append("\n")
                
                return "\n".join(content_parts)
            finally:
                os.unlink(temp_file_path)
                
        except Exception as e:
            print(f"Error extracting Excel content: {e}")
            return ""
    
    def redact_content(self, file_path: str, replacements: Dict[str, str]) -> str:
        """Redact Excel content across all sheets"""
        if not replacements:
            return file_path
        
        try:
            # Download original file
            file_response = w.files.download(file_path=file_path)
            
            with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as temp_file:
                if hasattr(file_response, 'content'):
                    temp_file.write(file_response.content)
                elif hasattr(file_response, 'read'):
                    temp_file.write(file_response.read())
                else:
                    temp_file.write(file_response)
                temp_file_path = temp_file.name
            
            try:
                # Load and process workbook
                workbook = openpyxl.load_workbook(temp_file_path)
                
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    
                    for row in sheet.iter_rows():
                        for cell in row:
                            if cell.value and isinstance(cell.value, str):
                                original_value = cell.value
                                for search_text, replace_text in replacements.items():
                                    cell.value = re.sub(re.escape(search_text), replace_text, cell.value, flags=re.IGNORECASE)
                
                # Save redacted file
                redacted_path = generate_output_filename(file_path, "redacted", ".xlsx")
                
                with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as output_temp:
                    workbook.save(output_temp.name)
                    
                    with open(output_temp.name, 'rb') as f:
                        w.files.upload(
                            file_path=redacted_path,
                            contents=f,
                            overwrite=True
                        )
                    os.unlink(output_temp.name)
                
                return redacted_path
            finally:
                os.unlink(temp_file_path)
                
        except Exception as e:
            print(f"Error redacting Excel content: {e}")
            raise e
    
    def get_supported_extensions(self) -> List[str]:
        return ['.xlsx', '.xls']

# PowerPoint handler implementation  
class PowerPointHandler(BaseFormatHandler):
    def can_handle(self, file_path: str, mime_type: str = None) -> bool:
        return (file_path.lower().endswith(('.pptx', '.ppt')) or 
                mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation') and POWERPOINT_SUPPORT
    
    def extract_content(self, file_path: str) -> str:
        """Extract text content from PowerPoint slides"""
        if not w:
            raise Exception("Databricks connection not configured")
        
        try:
            # Download file from UC Volume
            file_response = w.files.download(file_path=file_path)
            
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_file:
                if hasattr(file_response, 'content'):
                    temp_file.write(file_response.content)
                elif hasattr(file_response, 'read'):
                    temp_file.write(file_response.read())
                else:
                    temp_file.write(file_response)
                temp_file_path = temp_file.name
            
            try:
                # Read PowerPoint file
                presentation = Presentation(temp_file_path)
                content_parts = []
                
                for i, slide in enumerate(presentation.slides):
                    content_parts.append(f"## Slide {i+1}\n")
                    
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            content_parts.append(shape.text)
                    content_parts.append("\n")
                
                return "\n".join(content_parts)
            finally:
                os.unlink(temp_file_path)
                
        except Exception as e:
            print(f"Error extracting PowerPoint content: {e}")
            return ""
    
    def redact_content(self, file_path: str, replacements: Dict[str, str]) -> str:
        """Redact PowerPoint content across all slides"""
        if not replacements:
            return file_path
        
        try:
            # Download original file
            file_response = w.files.download(file_path=file_path)
            
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_file:
                if hasattr(file_response, 'content'):
                    temp_file.write(file_response.content)
                elif hasattr(file_response, 'read'):
                    temp_file.write(file_response.read())
                else:
                    temp_file.write(file_response)
                temp_file_path = temp_file.name
            
            try:
                # Load and process presentation
                presentation = Presentation(temp_file_path)
                
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            original_text = shape.text
                            for search_text, replace_text in replacements.items():
                                shape.text = re.sub(re.escape(search_text), replace_text, shape.text, flags=re.IGNORECASE)
                
                # Save redacted file
                redacted_path = generate_output_filename(file_path, "redacted", ".pptx")
                
                with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as output_temp:
                    presentation.save(output_temp.name)
                    
                    with open(output_temp.name, 'rb') as f:
                        w.files.upload(
                            file_path=redacted_path,
                            contents=f,
                            overwrite=True
                        )
                    os.unlink(output_temp.name)
                
                return redacted_path
            finally:
                os.unlink(temp_file_path)
                
        except Exception as e:
            print(f"Error redacting PowerPoint content: {e}")
            raise e
    
    def get_supported_extensions(self) -> List[str]:
        return ['.pptx', '.ppt']

# Format handler registry
FORMAT_HANDLERS = {
    'markdown': MarkdownHandler(),
    'excel': ExcelHandler(),
    'powerpoint': PowerPointHandler()
}

def get_format_handler(format_type: str) -> Optional[BaseFormatHandler]:
    """Get appropriate format handler"""
    return FORMAT_HANDLERS.get(format_type)

# Base exporter class
class BaseExporter(ABC):
    @abstractmethod
    def export(self, content: str, metadata: Dict[str, Any]) -> str:
        """Export content to specific format, return file path"""
        pass
    
    @abstractmethod
    def get_mime_type(self) -> str:
        """Return MIME type for exported format"""
        pass

# Markdown exporter implementation
class MarkdownExporter(BaseExporter):
    def export(self, content: str, metadata: Dict[str, Any]) -> str:
        """Export content as Markdown file"""
        try:
            # Format content as proper Markdown
            output_content = f"# Exported Document\n\n"
            
            if metadata.get("files"):
                output_content += f"**Source Files:** {', '.join(metadata['files'])}\n\n"
            
            import datetime
            output_content += f"**Export Date:** {datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n"
            output_content += "---\n\n"
            output_content += content
            
            # Create temporary file and upload to UC Volume
            timestamp = datetime.datetime.now().strftime('%Y%m%d_%H%M%S')
            export_filename = f"exported_document_{timestamp}.md"
            export_path = f"{get_uc_volume_path().rstrip('/')}/{export_filename}"
            
            with tempfile.NamedTemporaryFile(mode='w', encoding='utf-8', delete=False, suffix='.md') as temp_file:
                temp_file.write(output_content)
                temp_file_path = temp_file.name
            
            try:
                with open(temp_file_path, 'rb') as f:
                    w.files.upload(
                        file_path=export_path,
                        contents=f,
                        overwrite=True
                    )
            finally:
                os.unlink(temp_file_path)
            
            return export_path
        except Exception as e:
            print(f"Error exporting to Markdown: {e}")
            raise e
    
    def get_mime_type(self) -> str:
        return 'text/markdown'

# Excel exporter implementation
class ExcelExporter(BaseExporter):
    def export(self, content: str, metadata: Dict[str, Any]) -> str:
        """Export content as Excel workbook"""
        if not EXCEL_SUPPORT:
            raise Exception("Excel support not available")
        
        try:
            # Create workbook
            workbook = Workbook()
            
            # Remove default sheet and create new ones
            workbook.remove(workbook.active)
            
            # Create summary sheet
            summary_sheet = workbook.create_sheet("Summary")
            summary_sheet['A1'] = "Exported Document Summary"
            from openpyxl.styles import Font
            summary_sheet['A1'].font = Font(bold=True, size=14)
            
            if metadata.get("files"):
                summary_sheet['A3'] = "Source Files:"
                summary_sheet['A3'].font = Font(bold=True)
                for i, file_path in enumerate(metadata['files']):
                    summary_sheet[f'A{4+i}'] = file_path
            
            # Create content sheet
            content_sheet = workbook.create_sheet("Content")
            content_sheet['A1'] = "Document Content"
            content_sheet['A1'].font = Font(bold=True, size=14)
            
            # Split content into lines and add to sheet
            lines = content.split('\n')
            for i, line in enumerate(lines):
                content_sheet[f'A{i+3}'] = line
            
            # Save and upload
            import datetime
            timestamp = datetime.datetime.now().strftime('%Y%m%d_%H%M%S')
            export_filename = f"exported_document_{timestamp}.xlsx"
            export_path = f"{get_uc_volume_path().rstrip('/')}/{export_filename}"
            
            with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as temp_file:
                workbook.save(temp_file.name)
                
                with open(temp_file.name, 'rb') as f:
                    w.files.upload(
                        file_path=export_path,
                        contents=f,
                        overwrite=True
                    )
                os.unlink(temp_file.name)
            
            return export_path
        except Exception as e:
            print(f"Error exporting to Excel: {e}")
            raise e
    
    def get_mime_type(self) -> str:
        return 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'

# PowerPoint exporter implementation
class PowerPointExporter(BaseExporter):
    def export(self, content: str, metadata: Dict[str, Any]) -> str:
        """Export content as PowerPoint presentation"""
        if not POWERPOINT_SUPPORT:
            raise Exception("PowerPoint support not available")
        
        try:
            # Create presentation
            presentation = Presentation()
            
            # Title slide
            title_slide_layout = presentation.slide_layouts[0]
            title_slide = presentation.slides.add_slide(title_slide_layout)
            title = title_slide.shapes.title
            subtitle = title_slide.placeholders[1]
            
            title.text = "Exported Document"
            subtitle.text = f"Generated from {len(metadata.get('files', []))} source file(s)"
            
            # Content slides - split content into sections
            lines = content.split('\n')
            current_slide = None
            slide_layout = presentation.slide_layouts[1]  # Title and Content layout
            
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                
                # Check if this should be a new slide (starts with ##)
                if line.startswith('##'):
                    current_slide = presentation.slides.add_slide(slide_layout)
                    title_shape = current_slide.shapes.title
                    title_shape.text = line.replace('##', '').strip()
                    
                    # Add content placeholder
                    content_shape = current_slide.placeholders[1]
                    content_shape.text = ""
                elif current_slide is not None:
                    # Add content to current slide
                    content_shape = current_slide.placeholders[1]
                    if content_shape.text:
                        content_shape.text += f"\n{line}"
                    else:
                        content_shape.text = line
                else:
                    # Create first content slide if none exists
                    current_slide = presentation.slides.add_slide(slide_layout)
                    title_shape = current_slide.shapes.title
                    title_shape.text = "Document Content"
                    content_shape = current_slide.placeholders[1]
                    content_shape.text = line
            
            # Save and upload
            import datetime
            timestamp = datetime.datetime.now().strftime('%Y%m%d_%H%M%S')
            export_filename = f"exported_document_{timestamp}.pptx"
            export_path = f"{get_uc_volume_path().rstrip('/')}/{export_filename}"
            
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_file:
                presentation.save(temp_file.name)
                
                with open(temp_file.name, 'rb') as f:
                    w.files.upload(
                        file_path=export_path,
                        contents=f,
                        overwrite=True
                    )
                os.unlink(temp_file.name)
            
            return export_path
        except Exception as e:
            print(f"Error exporting to PowerPoint: {e}")
            raise e
    
    def get_mime_type(self) -> str:
        return 'application/vnd.openxmlformats-officedocument.presentationml.presentation'

# Exporter registry
EXPORTERS = {
    'md': MarkdownExporter(),
    'xlsx': ExcelExporter(),
    'pptx': PowerPointExporter()
}

def get_exporter(export_format: str) -> Optional[BaseExporter]:
    """Get appropriate exporter"""
    return EXPORTERS.get(export_format)

class WarehouseConfigRequest(BaseModel):
    warehouse_id: str

class VolumePathConfigRequest(BaseModel):
    volume_path: str

class DeltaTablePathConfigRequest(BaseModel):
    delta_table_path: str


class ParseDocumentRequest(BaseModel):
    file_path: str


class RedactPDFRequest(BaseModel):
    file_paths: List[str]

# New Pydantic models for multi-format support
class ExportRequest(BaseModel):
    file_paths: List[str]
    export_format: str  # 'md', 'xlsx', 'pptx'
    output_filename: Optional[str] = None

class RedactDocumentRequest(BaseModel):
    file_paths: List[str]
    file_types: List[str] = []  # Support multiple formats

class FileFormatInfo(BaseModel):
    supported_formats: List[str]
    upload_formats: List[str]
    export_formats: List[str]
    format_handlers_available: Dict[str, bool]


@app.get("/api/warehouse-config")
def get_warehouse_config():
    """Get current warehouse configuration"""
    return {
        "warehouse_id": current_warehouse_id,
        "default_warehouse_id": warehouse_id
    }

@app.post("/api/warehouse-config")
def update_warehouse_config(request: WarehouseConfigRequest):
    """Update warehouse configuration"""
    global current_warehouse_id
    current_warehouse_id = request.warehouse_id
    print(f"🔧 Warehouse ID updated to: {current_warehouse_id}")
    return {
        "success": True,
        "warehouse_id": current_warehouse_id,
        "message": "Warehouse ID updated successfully"
    }

@app.get("/api/volume-path-config")
def get_volume_path_config():
    """Get current volume path configuration"""
    default_path = YAML_CONFIG.get("DATABRICKS_VOLUME_PATH", "/Volumes/fins_genai/unstructured_documents/pdf_tpg/")
    return {
        "volume_path": current_volume_path or default_path,
        "default_volume_path": default_path
    }

@app.post("/api/volume-path-config")
def update_volume_path_config(request: VolumePathConfigRequest):
    """Update volume path configuration"""
    global current_volume_path
    current_volume_path = request.volume_path
    print(f"🔧 Volume path updated to: {current_volume_path}")
    return {
        "success": True,
        "volume_path": current_volume_path,
        "message": "Volume path updated successfully"
    }

@app.get("/api/delta-table-path-config")
def get_delta_table_path_config():
    """Get current delta table path configuration"""
    default_path = YAML_CONFIG.get("DATABRICKS_DELTA_TABLE_PATH", "/fins_genai.unstructured_documents.files_parsed")
    return {
        "delta_table_path": current_delta_table_path or default_path,
        "default_delta_table_path": default_path
    }

@app.post("/api/delta-table-path-config")
def update_delta_table_path_config(request: DeltaTablePathConfigRequest):
    """Update delta table path configuration"""
    global current_delta_table_path
    current_delta_table_path = request.delta_table_path
    print(f"🔧 Delta table path updated to: {current_delta_table_path}")
    return {
        "success": True,
        "delta_table_path": current_delta_table_path,
        "message": "Delta table path updated successfully"
    }


@app.post("/api/upload-to-uc")
async def upload_to_uc(files: List[UploadFile] = FastAPIFile(...)):
    """Upload files to Databricks UC Volume - supports multiple formats"""
    if not w:
        raise HTTPException(status_code=500, detail="Databricks connection is not configured.")
    
    try:
        uploaded_files = []
        
        for file in files:
            # Validate file format
            file_ext = get_file_extension(file.filename)
            if file_ext not in SUPPORTED_UPLOAD_FORMATS:
                raise HTTPException(
                    status_code=400, 
                    detail=f"Unsupported file format: {file_ext}. Supported formats: {', '.join(SUPPORTED_UPLOAD_FORMATS)}"
                )
            # Create a temporary file to store the uploaded content
            with tempfile.NamedTemporaryFile(delete=False) as temp_file:
                # Copy file content to temporary file
                shutil.copyfileobj(file.file, temp_file)
                temp_file_path = temp_file.name
            
            try:
                # Upload to UC Volume - fix double slash issue
                base_path = get_uc_volume_path().rstrip('/')  # Remove trailing slash
                uc_file_path = f"{base_path}/{file.filename}"
                
                # Upload to UC Volume using the Files API with file handle
                with open(temp_file_path, 'rb') as f:
                    w.files.upload(
                        file_path=uc_file_path,
                        contents=f,
                        overwrite=True
                    )
                
                # Get file size for response
                file_size = os.path.getsize(temp_file_path)
                
                # Detect format and handler availability
                file_format = detect_file_format(file.filename)
                handler_available = get_format_handler(file_format) is not None or file_format == 'pdf'
                
                uploaded_files.append({
                    "name": file.filename,
                    "path": uc_file_path,
                    "size": file_size,
                    "format": file_format,
                    "extension": file_ext,
                    "handler_available": handler_available,
                    "redaction_supported": handler_available
                })
                
            finally:
                # Clean up temporary file
                os.unlink(temp_file_path)
        
        return {
            "success": True,
            "uploaded_files": uploaded_files,
            "message": f"Successfully uploaded {len(uploaded_files)} files to UC Volume"
        }
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Upload failed: {str(e)}")

@app.post("/api/test-ai-functions")
def test_ai_functions():
    """Test if AI Functions are available and working"""
    if not w:
        raise HTTPException(status_code=500, detail="Databricks connection is not configured.")
    
    if not current_warehouse_id:
        raise HTTPException(status_code=500, detail="DATABRICKS_WAREHOUSE_ID is not set.")
    
    try:
        # Test basic AI function availability with simple text
        test_query = """
        SELECT 
            ai_extract('This is a test document about John Doe from ACME Corp', 
                      ARRAY('person_name', 'company_name')) as extracted_info
        """
        
        print(f"Testing AI Functions availability: {test_query}")
        
        result = w.statement_execution.execute_statement(
            statement=test_query,
            warehouse_id=current_warehouse_id,
            wait_timeout='30s'
        )
        
        print(f"AI Functions test result status: {result.status}")
        
        if result.result and result.result.data_array:
            return {
                "success": True,
                "message": "AI Functions are working",
                "test_result": result.result.data_array[0][0] if result.result.data_array[0] else None
            }
        else:
            return {
                "success": False,
                "message": "AI Functions test returned no data"
            }
            
    except Exception as e:
        print(f"AI Functions test error: {e}")
        error_msg = str(e)
        
        if "FUNCTION_NOT_FOUND" in error_msg or "ai_extract" in error_msg:
            error_msg = "AI Functions not available - ensure they are enabled for your warehouse"
        elif "PERMISSION_DENIED" in error_msg:
            error_msg = "Permission denied - check warehouse permissions for AI Functions"
            
        return {
            "success": False,
            "message": error_msg,
            "error_type": type(e).__name__
        }

@app.post("/api/write-to-delta-table")
def write_to_delta_table(request: WriteToTableRequest):
    """Write processed documents to delta table using ai_parse_document - replaces entire table"""
    if not w:
        raise HTTPException(status_code=500, detail="Databricks connection is not configured.")
    
    if not current_warehouse_id:
        raise HTTPException(status_code=500, detail="DATABRICKS_WAREHOUSE_ID is not set.")

    if not request.file_paths:
        raise HTTPException(status_code=400, detail="file_paths is required")
    
    # Validate we have exactly one file
    if len(request.file_paths) != 1:
        raise HTTPException(status_code=400, detail="Only one file can be processed at a time")

    try:
        # Get the single file path
        file_path = request.file_paths[0]
        
        # Get the existing delta table path
        destination_table = get_delta_table_path()
        print(f"Working with delta table: {destination_table}")
        print(f"Processing single file: {file_path}")
        
        # Check if table exists and get its schema
        check_table_query = f"""
        DESCRIBE IDENTIFIER('{destination_table}')
        """
        
        print("Checking table schema...")
        try:
            # Check if table has new schema
            has_new_schema = False
            if not has_new_schema:
                print("Table has old schema or doesn't exist. Creating/recreating table...")
                
                # First drop the table
                drop_query = f"DROP TABLE IF EXISTS IDENTIFIER('{destination_table}')"
                
                drop_result = w.statement_execution.execute_statement(
                    statement=drop_query,
                    warehouse_id=current_warehouse_id,
                    wait_timeout='30s'
                )
                
                if drop_result.status and drop_result.status.state == StatementState.FAILED:
                    raise Exception(f"Failed to drop table: {drop_result.status}")
                
                # Then create the table with new schema
                create_query = f"""
                CREATE TABLE IDENTIFIER('{destination_table}') (
                    path STRING,
                    content STRING
                ) USING DELTA
                """
                
                create_result = w.statement_execution.execute_statement(
                    statement=create_query,
                    warehouse_id=current_warehouse_id,
                    wait_timeout='30s'
                )
                
                if create_result.status and create_result.status.state == StatementState.FAILED:
                    raise Exception(f"Failed to create table: {create_result.status}")
                    
                print("Table recreated with new schema")
            else:
                print("Table already has correct schema")
                
        except Exception as e:
            if "TABLE_OR_VIEW_NOT_FOUND" in str(e):
                print("Table doesn't exist, creating new table...")
                create_table_query = f"""
                CREATE TABLE IDENTIFIER('{destination_table}') (
                    path STRING,
                    content STRING
                ) USING DELTA
                """
                
                create_result = w.statement_execution.execute_statement(
                    statement=create_table_query,
                    warehouse_id=current_warehouse_id,
                    wait_timeout='30s'
                )
                
                if create_result.status and create_result.status.state == StatementState.FAILED:
                    raise Exception(f"Failed to create table: {create_result.status}")
            else:
                raise e
        
        print("Table exists with correct schema, replacing all data...")
        
        # Convert to dbfs format for the path column
        if file_path.startswith('/Volumes/'):
            dbfs_path = 'dbfs:' + file_path
        else:
            dbfs_path = file_path
        
        print(f"DBFS path will be: {dbfs_path}")
        
        # TRUNCATE entire table - delete ALL existing records
        truncate_query = f"""
        DELETE FROM IDENTIFIER('{destination_table}')
        """
        
        print(f"Truncating entire table...")
        truncate_result = w.statement_execution.execute_statement(
            statement=truncate_query,
            warehouse_id=current_warehouse_id,
            wait_timeout='30s'
        )
        
        if truncate_result.status and truncate_result.status.state == StatementState.FAILED:
            print(f"Truncate operation failed: {truncate_result.status}")
        else:
            print("Table truncated successfully")
        
        # Then insert new records from the single file with deterministic table IDs
        insert_query = f"""
        INSERT INTO IDENTIFIER('{destination_table}')
        WITH file_data AS (
          SELECT 
            path,
            content
          FROM READ_FILES('{dbfs_path}', format => 'binaryFile')
        ),
        parsed_documents AS (
          SELECT
              path,
              ai_parse_document(content) as parsed,
              content
          FROM file_data
        ),
        -- Extract page markdowns from ai_parse output
        sorted_page_contents AS (
          SELECT
            path,
            page:content AS content
          FROM
            (
              SELECT
                path,
                posexplode(try_cast(parsed:document:pages AS ARRAY<VARIANT>)) AS (page_idx, page)
              FROM
                parsed_documents
              WHERE
                parsed:document:pages IS NOT NULL
                AND CAST(parsed:error_status AS STRING) IS NULL
            )
          ORDER BY
            page_idx
        ),
        -- Concatenate so we have 1 row per document
        concatenated AS (
            SELECT
                path,
                concat_ws('
        
        ', collect_list(content)) AS full_content
            FROM
                sorted_page_contents
            GROUP BY
                path
        ),
        -- Bring back the raw parsing since it could be useful for other downstream uses
        parsed AS (
            SELECT
                a.path,
                b.parsed as raw_parsed,
                a.full_content as content
            FROM concatenated a
            JOIN parsed_documents b ON a.path = b.path
        )
        SELECT path, content FROM parsed
        """
        
        print(f"Executing INSERT for {file_path}")
        
        insert_result = w.statement_execution.execute_statement(
            statement=insert_query,
            warehouse_id=current_warehouse_id,
            wait_timeout='50s'
        )
        
        print(f"INSERT result: {insert_result.status}")
        
        # If the operation is still pending or running, wait for it to complete
        if insert_result.status and insert_result.status.state in [StatementState.PENDING, StatementState.RUNNING]:
            print(f"INSERT operation is pending, waiting for completion...")
            try:
                # Wait for the statement to complete
                final_result = w.statement_execution.get_statement(insert_result.statement_id)
                
                # Keep checking until it's no longer pending or running (up to additional 30 seconds)
                import time
                max_wait = 300
                waited = 0
                while final_result.status.state in [StatementState.PENDING, StatementState.RUNNING] and waited < max_wait:
                    time.sleep(2)
                    waited += 2
                    final_result = w.statement_execution.get_statement(insert_result.statement_id)
                    print(f"Waiting for INSERT completion... ({waited}s) - Status: {final_result.status.state}")
                
                print(f"Final INSERT result: {final_result.status}")
                insert_result = final_result
                
            except Exception as wait_error:
                print(f"Error waiting for INSERT completion: {wait_error}")
        
        if insert_result.status and insert_result.status.state == StatementState.SUCCEEDED:
            print(f"Successfully processed: {file_path} (read from: {dbfs_path})")
            
            return {
                "success": True,
                "destination_table": destination_table,
                "processed_files": [file_path],
                "message": f"Successfully extracted tables from document and replaced entire table"
            }
        else:
            error_msg = f"Failed to process {file_path}"
            if insert_result.status and insert_result.status.error:
                error_msg += f": {insert_result.status.error}"
            print(error_msg)
            
            return {
                "success": False,
                "destination_table": destination_table,
                "processed_files": [file_path],
                "processed_paths": [],
                "data": [],
                "total_results": 0,
                "message": error_msg
            }

    except Exception as e:
        print(f"Delta table write error: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to write to delta table: {str(e)}")

@app.post("/api/query-delta-table")
def query_delta_table(request: QueryDeltaTableRequest):
    """Query delta table results for specific documents"""
    if not w:
        raise HTTPException(status_code=500, detail="Databricks connection is not configured.")
    
    if not current_warehouse_id:
        raise HTTPException(status_code=500, detail="DATABRICKS_WAREHOUSE_ID is not set.")

    try:
        # Get the delta table path
        destination_table = get_delta_table_path()
        print(f"Querying delta table: {destination_table}")
        
        # Build the query with optional file filtering
        where_clause = ""
        if request.file_paths:
            # Convert to dbfs: format for filtering
            dbfs_file_paths = []
            for fp in request.file_paths:
                if fp.startswith('/Volumes/'):
                    dbfs_path = 'dbfs:' + fp
                else:
                    dbfs_path = fp
                dbfs_file_paths.append(dbfs_path)
            
            # Use exact path matching instead of LIKE with filename
            path_conditions = ", ".join([f"'{fp}'" for fp in dbfs_file_paths])
            where_clause = f"WHERE path IN ({path_conditions})"
        
        query = f"""
        SELECT
            path,
            content
        FROM IDENTIFIER('{destination_table}')
        {where_clause}
        LIMIT {request.limit}
        """
        
        print(f"Executing query: {query}")
        
        result = w.statement_execution.execute_statement(
            statement=query,
            warehouse_id=current_warehouse_id,
            wait_timeout='30s'
        )

        if result.result and result.result.data_array:
            delta_results = []
            for row in result.result.data_array:
                delta_results.append({
                    "path": row[0] if len(row) > 0 else "",
                    "content": row[1] if len(row) > 1 else "",
                })
            
            print(f"Returning {len(delta_results)} results from delta table")
            return {
                "success": True,
                "data": delta_results,
                "table_name": destination_table,
                "total_results": len(delta_results)
            }
        else:
            print("No data returned from query")
            return {
                "success": True,
                "data": [],
                "message": "No results found in delta table"
            }

    except Exception as e:
        print(f"Delta table query error: {e}")
        return {
            "success": False,
            "data": [],
            "error": f"Failed to query delta table: {str(e)}"
        }


@app.get("/api/download-redacted-pdf")
def download_redacted_pdf(file_path: str):
    """Download a redacted PDF file from UC Volume"""
    if not w:
        raise HTTPException(status_code=500, detail="Databricks connection is not configured.")
    
    if not file_path:
        raise HTTPException(status_code=400, detail="file_path parameter is required")
    
    try:
        print(f"Downloading redacted PDF file: {file_path}")
        
        # Download the PDF file from UC Volume
        file_response = w.files.download(file_path=file_path)
        
        # Create a temporary file to store the downloaded content
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
            temp_file_path = temp_file.name
            
            # Handle different response types from Databricks SDK (same as redaction function)
            response_type = type(file_response).__name__
            print(f"Download response type: {response_type}")
            
            try:
                # Use the same robust handling as the redaction function
                if hasattr(file_response, 'iter_content') and callable(getattr(file_response, 'iter_content')):
                    print("Using iter_content method")
                    for chunk in file_response.iter_content(chunk_size=8192):
                        if chunk:
                            temp_file.write(chunk)
                elif hasattr(file_response, 'content'):
                    print("Using content attribute")
                    content_data = getattr(file_response, 'content')
                    if isinstance(content_data, bytes):
                        temp_file.write(content_data)
                    else:
                        raise Exception(f"Content is not bytes: {type(content_data)}")
                elif hasattr(file_response, 'contents'):
                    print("Using contents attribute")
                    contents = getattr(file_response, 'contents')
                    if isinstance(contents, bytes):
                        temp_file.write(contents)
                    else:
                        # Handle StreamingResponse within contents
                        if hasattr(contents, 'iter_content') and callable(getattr(contents, 'iter_content')):
                            print("Contents has iter_content method")
                            for chunk in contents.iter_content(chunk_size=8192):
                                if chunk:
                                    temp_file.write(chunk)
                        elif hasattr(contents, 'content'):
                            print("Contents has content attribute")
                            inner_content = getattr(contents, 'content')
                            if isinstance(inner_content, bytes):
                                temp_file.write(inner_content)
                            else:
                                raise Exception(f"Inner content is not bytes: {type(inner_content)}")
                        elif hasattr(contents, 'read') and callable(getattr(contents, 'read')):
                            print("Contents has read method")
                            read_content = contents.read()
                            if isinstance(read_content, bytes):
                                temp_file.write(read_content)
                            else:
                                temp_file.write(read_content.encode() if isinstance(read_content, str) else read_content)
                        else:
                            raise Exception(f"Cannot handle StreamingResponse contents: {type(contents)}")
                elif hasattr(file_response, 'read') and callable(getattr(file_response, 'read')):
                    print("Using read method")
                    content_data = file_response.read()
                    if isinstance(content_data, bytes):
                        temp_file.write(content_data)
                    else:
                        temp_file.write(content_data.encode() if isinstance(content_data, str) else content_data)
                elif isinstance(file_response, bytes):
                    print("Direct bytes response")
                    temp_file.write(file_response)
                else:
                    raise Exception(f"Unsupported download response type: {response_type}. Available attributes: {[attr for attr in dir(file_response) if not attr.startswith('_')][:10]}")
                    
            except Exception as e:
                print(f"Error writing to temp file: {e}")
                os.unlink(temp_file_path)  # Clean up temp file
                raise Exception(f"Failed to process download response: {str(e)}")
        
        # Check if temp file has content
        file_size = os.path.getsize(temp_file_path)
        print(f"Temp file size: {file_size} bytes")
        
        if file_size == 0:
            os.unlink(temp_file_path)  # Clean up empty temp file
            raise HTTPException(status_code=404, detail="Redacted PDF file not found or empty")
        
        # Extract filename from path
        filename = os.path.basename(file_path)
        if not filename.endswith('.pdf'):
            filename += '.pdf'
        
        # Return the file as a download using FileResponse
        def cleanup_temp_file():
            """Clean up temp file after response is sent"""
            try:
                os.unlink(temp_file_path)
                print(f"Cleaned up temp file: {temp_file_path}")
            except:
                pass
        
        # Use FileResponse for better file handling
        response = FileResponse(
            path=temp_file_path,
            media_type="application/pdf",
            filename=filename
        )
        
        # Schedule cleanup
        atexit.register(cleanup_temp_file)
        
        return response
        
    except Exception as e:
        print(f"Redacted PDF download error: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to download redacted PDF file: {str(e)}")


def extract_entities_for_redaction(content: str) -> dict:
    """
    Use Databricks AI to identify entities that need redaction using NER prompt
    """
    if not w or not current_warehouse_id:
        raise Exception("Databricks connection or warehouse not configured")
    
    # Load the NER prompt
    ner_prompt_path = os.path.join(os.path.dirname(__file__), "ner_prompt.md")
    try:
        with open(ner_prompt_path, 'r') as f:
            ner_prompt = f.read()
    except FileNotFoundError:
        raise Exception("NER prompt file not found")
    
    # Construct the AI query with the content
    full_prompt = f"""{ner_prompt}

## DOCUMENT TO ANALYZE

{content}

Please analyze the above document content and return ONLY a JSON dictionary of entities that need redaction, following the format specified in the prompt."""

    try:
        # Use ai_query function with the Claude Sonnet endpoint
        query = f"""
        SELECT ai_query(
            'databricks-claude-sonnet-4',
            '{full_prompt.replace("'", "''")}',
            'user'
        ) as ner_result
        """
        
        print(f"Executing NER query...")
        
        result = w.statement_execution.execute_statement(
            statement=query,
            warehouse_id=current_warehouse_id,
            wait_timeout='50s'  # Maximum allowed timeout for ai_query
        )
        
        if result.result and result.result.data_array and len(result.result.data_array) > 0:
            ner_response = result.result.data_array[0][0]
            print(f"Raw NER response: {ner_response}")
            
            # Try to extract JSON from the response
            try:
                # Look for JSON in the response
                import re
                json_match = re.search(r'\{.*\}', ner_response, re.DOTALL)
                if json_match:
                    json_str = json_match.group(0)
                    entities_dict = json.loads(json_str)
                    print(f"Extracted {len(entities_dict)} entities for redaction")
                    return entities_dict
                else:
                    print("No JSON found in NER response")
                    return {}
            except json.JSONDecodeError as e:
                print(f"Failed to parse JSON from NER response: {e}")
                return {}
        else:
            print("No result from NER query")
            return {}
            
    except Exception as e:
        print(f"NER query error: {e}")
        raise Exception(f"Failed to extract entities: {str(e)}")

def redact_pdf_from_uc(file_path: str, replacements: dict) -> str:
    """
    Download PDF from UC Volume, redact it, and upload back to UC Volume
    """
    if not w:
        raise Exception("Databricks connection not configured")
    
    if not replacements:
        print("No entities to redact, skipping redaction")
        return file_path
    
    # Create temporary files
    temp_input_path = None
    temp_output_path = None
    
    try:
        # Download the PDF from UC Volume
        print(f"Downloading PDF from UC Volume: {file_path}")
        file_response = w.files.download(file_path=file_path)
        
        # Create temporary input file
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_input:
            temp_input_path = temp_input.name
            
            # Handle different response types from Databricks SDK
            response_type = type(file_response).__name__
            print(f"File response type: {response_type}")
            
            # Write downloaded content to temp file
            if hasattr(file_response, 'iter_content') and callable(getattr(file_response, 'iter_content')):
                print("Using iter_content method")
                for chunk in file_response.iter_content(chunk_size=8192):
                    if chunk:
                        temp_input.write(chunk)
            elif hasattr(file_response, 'content'):
                print("Using content attribute")
                content_data = getattr(file_response, 'content')
                if isinstance(content_data, bytes):
                    temp_input.write(content_data)
                else:
                    raise Exception(f"Content is not bytes: {type(content_data)}")
            elif hasattr(file_response, 'contents'):
                print("Using contents attribute")
                contents = getattr(file_response, 'contents')
                if isinstance(contents, bytes):
                    temp_input.write(contents)
                else:
                    # Handle StreamingResponse within contents
                    if hasattr(contents, 'iter_content') and callable(getattr(contents, 'iter_content')):
                        print("Contents has iter_content method")
                        for chunk in contents.iter_content(chunk_size=8192):
                            if chunk:
                                temp_input.write(chunk)
                    elif hasattr(contents, 'content'):
                        print("Contents has content attribute")
                        inner_content = getattr(contents, 'content')
                        if isinstance(inner_content, bytes):
                            temp_input.write(inner_content)
                        else:
                            raise Exception(f"Inner content is not bytes: {type(inner_content)}")
                    elif hasattr(contents, 'read') and callable(getattr(contents, 'read')):
                        print("Contents has read method")
                        read_content = contents.read()
                        if isinstance(read_content, bytes):
                            temp_input.write(read_content)
                        else:
                            temp_input.write(read_content.encode() if isinstance(read_content, str) else read_content)
                    else:
                        raise Exception(f"Cannot handle StreamingResponse contents: {type(contents)}")
            elif hasattr(file_response, 'read') and callable(getattr(file_response, 'read')):
                print("Using read method")
                content_data = file_response.read()
                if isinstance(content_data, bytes):
                    temp_input.write(content_data)
                else:
                    temp_input.write(content_data.encode() if isinstance(content_data, str) else content_data)
            elif isinstance(file_response, bytes):
                print("Direct bytes response")
                temp_input.write(file_response)
            else:
                raise Exception(f"Unsupported file response type: {response_type}. Available attributes: {[attr for attr in dir(file_response) if not attr.startswith('_')][:10]}")
        
        # Perform redaction using PyMuPDF
        print(f"Performing redaction with {len(replacements)} entities")
        doc = fitz.open(temp_input_path)
        
        if doc.is_encrypted:
            doc.close()
            raise Exception("PDF is password protected")
        
        changes_made = False
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            
            for search_text, replace_text in replacements.items():
                # Search for text instances (case-insensitive)
                areas = page.search_for(search_text)
                
                if areas:
                    changes_made = True
                    print(f"Found '{search_text}' on page {page_num + 1}, replacing with '{replace_text}'")
                    for rect in areas:
                        # Add redaction annotation with white fill and replacement text
                        page.add_redact_annot(rect, text=replace_text, 
                                            fill=(1, 1, 1), text_color=(0, 0, 0))
            
            # Apply redactions for this page
            page.apply_redactions()
        
        # Save to temporary output file
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_output:
            temp_output_path = temp_output.name
        
        if changes_made:
            doc.save(temp_output_path, garbage=4, deflate=True)
            print(f"Redaction completed with {len(replacements)} entity types processed")
        else:
            # If no changes, just copy the original
            doc.save(temp_output_path)
            print("No matching text found for redaction")
        
        doc.close()
        
        # Generate redacted filename
        file_dir = os.path.dirname(file_path)
        file_name = os.path.basename(file_path)
        redacted_filename = f"redacted_{file_name}"
        redacted_path = os.path.join(file_dir, redacted_filename).replace('\\', '/')
        
        # Upload redacted PDF back to UC Volume
        print(f"Uploading redacted PDF to UC Volume: {redacted_path}")
        with open(temp_output_path, 'rb') as f:
            file_data = f.read()
            # Convert to BytesIO for upload
            binary_data = io.BytesIO(file_data)
            w.files.upload(
                file_path=redacted_path,
                contents=binary_data,
                overwrite=True
            )
        
        print(f"Successfully created redacted PDF: {redacted_path}")
        return redacted_path
        
    except Exception as e:
        print(f"PDF redaction error: {e}")
        raise e
    finally:
        # Clean up temporary files
        if temp_input_path and os.path.exists(temp_input_path):
            os.unlink(temp_input_path)
        if temp_output_path and os.path.exists(temp_output_path):
            os.unlink(temp_output_path)

@app.post("/api/redact-pdf")
def redact_pdf_documents(request: RedactPDFRequest):
    """
    Perform NER-based redaction on uploaded PDF documents
    """
    if not w:
        raise HTTPException(status_code=500, detail="Databricks connection is not configured.")
    
    if not current_warehouse_id:
        raise HTTPException(status_code=500, detail="DATABRICKS_WAREHOUSE_ID is not set.")
    
    if not request.file_paths:
        raise HTTPException(status_code=400, detail="file_paths is required")
    
    try:
        destination_table = get_delta_table_path()
        redacted_files = []
        
        for original_file_path in request.file_paths:
            print(f"Processing redaction for: {original_file_path}")
            
            # Check if this is a PDF file
            if not original_file_path.lower().endswith('.pdf'):
                print(f"Skipping non-PDF file: {original_file_path}")
                continue
            
            # Convert to dbfs format for querying
            if original_file_path.startswith('/Volumes/'):
                dbfs_path = 'dbfs:' + original_file_path
            else:
                dbfs_path = original_file_path
            
            # Get document content from Delta table
            query = f"""
            SELECT content
            FROM IDENTIFIER('{destination_table}')
            WHERE path = '{dbfs_path}'
            LIMIT 1
            """
            
            print(f"Querying content for: {original_file_path}")
            result = w.statement_execution.execute_statement(
                statement=query,
                warehouse_id=current_warehouse_id,
                wait_timeout='30s'
            )
            
            if not result.result or not result.result.data_array:
                print(f"No content found in Delta table for: {original_file_path}")
                continue
            
            content = result.result.data_array[0][0]
            if not content or not content.strip():
                print(f"Empty content found for: {original_file_path}")
                continue
            
            # Extract entities for redaction using NER
            print("Extracting entities using NER...")
            entities_to_redact = extract_entities_for_redaction(content)
            
            if not entities_to_redact:
                print("No entities identified for redaction")
                redacted_files.append({
                    "original_file": original_file_path,
                    "redacted_file": original_file_path,  # No redaction needed
                    "entities_count": 0,
                    "status": "no_entities_found"
                })
                continue
            
            # Perform PDF redaction
            print(f"Performing PDF redaction with {len(entities_to_redact)} entities")
            redacted_path = redact_pdf_from_uc(original_file_path, entities_to_redact)
            
            redacted_files.append({
                "original_file": original_file_path,
                "redacted_file": redacted_path,
                "entities_count": len(entities_to_redact),
                "entities": entities_to_redact,
                "status": "redacted"
            })
            
            print(f"Successfully processed: {original_file_path}")
        
        if not redacted_files:
            return {
                "success": False,
                "message": "No PDF files were processed for redaction",
                "redacted_files": []
            }
        
        return {
            "success": True,
            "message": f"Successfully processed {len(redacted_files)} file(s) for redaction",
            "redacted_files": redacted_files
        }
        
    except Exception as e:
        print(f"PDF redaction error: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to redact PDF files: {str(e)}")

# New endpoints for multi-format support

@app.get("/api/supported-formats")
def get_supported_formats():
    """Get list of supported file formats for upload and export"""
    return FileFormatInfo(
        supported_formats=SUPPORTED_UPLOAD_FORMATS,
        upload_formats=SUPPORTED_UPLOAD_FORMATS,
        export_formats=SUPPORTED_EXPORT_FORMATS,
        format_handlers_available={
            'markdown': MARKDOWN_SUPPORT,
            'excel': EXCEL_SUPPORT, 
            'powerpoint': POWERPOINT_SUPPORT,
            'pdf': True  # Always available
        }
    )

@app.post("/api/export-document")
def export_document(request: ExportRequest):
    """Export processed document data to specified format"""
    if not w:
        raise HTTPException(status_code=500, detail="Databricks connection is not configured.")
    
    if request.export_format not in EXPORTERS:
        raise HTTPException(
            status_code=400, 
            detail=f"Unsupported export format: {request.export_format}. Supported: {list(EXPORTERS.keys())}"
        )
    
    try:
        destination_table = get_delta_table_path()
        
        # Get document content from Delta table
        if request.file_paths:
            dbfs_paths = [('dbfs:' + fp if fp.startswith('/Volumes/') else fp) for fp in request.file_paths]
            path_conditions = ", ".join([f"'{fp}'" for fp in dbfs_paths])
            where_clause = f"WHERE path IN ({path_conditions})"
        else:
            where_clause = ""
        
        query = f"""
        SELECT path, content
        FROM IDENTIFIER('{destination_table}')
        {where_clause}
        """
        
        result = w.statement_execution.execute_statement(
            statement=query,
            warehouse_id=current_warehouse_id,
            wait_timeout='30s'
        )
        
        if not result.result or not result.result.data_array:
            raise HTTPException(status_code=404, detail="No document data found")
        
        # Combine content from all documents
        combined_content = ""
        metadata = {"files": [], "export_format": request.export_format}
        
        for row in result.result.data_array:
            path, content = row[0], row[1]
            combined_content += f"\n\n{content}" if combined_content else content
            metadata["files"].append(path)
        
        # Export using appropriate exporter
        exporter = EXPORTERS[request.export_format]
        export_file_path = exporter.export(combined_content, metadata)
        
        # Generate download filename
        output_filename = request.output_filename or f"exported_document.{request.export_format}"
        
        return {
            "success": True,
            "export_file_path": export_file_path,
            "download_filename": output_filename,
            "export_format": request.export_format,
            "files_processed": len(metadata["files"])
        }
        
    except Exception as e:
        print(f"Export error: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to export document: {str(e)}")

@app.post("/api/redact-document")
def redact_documents(request: RedactDocumentRequest):
    """Perform NER-based redaction on uploaded documents across formats"""
    if not w:
        raise HTTPException(status_code=500, detail="Databricks connection is not configured.")
    
    if not current_warehouse_id:
        raise HTTPException(status_code=500, detail="DATABRICKS_WAREHOUSE_ID is not set.")
    
    if not request.file_paths:
        raise HTTPException(status_code=400, detail="file_paths is required")
    
    try:
        destination_table = get_delta_table_path()
        redacted_files = []
        
        for original_file_path in request.file_paths:
            print(f"Processing redaction for: {original_file_path}")
            
            # Detect file format
            file_format = detect_file_format(original_file_path)
            handler = get_format_handler(file_format)
            
            # Fall back to PDF handler for PDF files
            if not handler and file_format == 'pdf':
                # Use existing PDF redaction logic
                if not original_file_path.lower().endswith('.pdf'):
                    print(f"Skipping non-PDF file: {original_file_path}")
                    continue
                
                # Convert to dbfs format for querying
                dbfs_path = 'dbfs:' + original_file_path if original_file_path.startswith('/Volumes/') else original_file_path
                
                # Get document content from Delta table
                query = f"""
                SELECT content
                FROM IDENTIFIER('{destination_table}')
                WHERE path = '{dbfs_path}'
                LIMIT 1
                """
                
                result = w.statement_execution.execute_statement(
                    statement=query,
                    warehouse_id=current_warehouse_id,
                    wait_timeout='30s'
                )
                
                if not result.result or not result.result.data_array:
                    continue
                
                content = result.result.data_array[0][0]
                if not content or not content.strip():
                    continue
                
                # Extract entities for redaction using NER
                entities_to_redact = extract_entities_for_redaction(content)
                
                if not entities_to_redact:
                    redacted_files.append({
                        "original_file": original_file_path,
                        "redacted_file": original_file_path,
                        "entities_count": 0,
                        "status": "no_entities_found",
                        "format": file_format
                    })
                    continue
                
                # Perform PDF redaction using existing function
                redacted_path = redact_pdf_from_uc(original_file_path, entities_to_redact)
                
                redacted_files.append({
                    "original_file": original_file_path,
                    "redacted_file": redacted_path,
                    "entities_count": len(entities_to_redact),
                    "entities": entities_to_redact,
                    "status": "redacted",
                    "format": file_format
                })
                continue
            
            if not handler:
                print(f"No handler available for format: {file_format} ({original_file_path})")
                redacted_files.append({
                    "original_file": original_file_path,
                    "redacted_file": original_file_path,
                    "entities_count": 0,
                    "status": "format_not_supported",
                    "format": file_format
                })
                continue
            
            # Get document content from Delta table
            dbfs_path = 'dbfs:' + original_file_path if original_file_path.startswith('/Volumes/') else original_file_path
            
            query = f"""
            SELECT content
            FROM IDENTIFIER('{destination_table}')
            WHERE path = '{dbfs_path}'
            LIMIT 1
            """
            
            result = w.statement_execution.execute_statement(
                statement=query,
                warehouse_id=current_warehouse_id,
                wait_timeout='30s'
            )
            
            if not result.result or not result.result.data_array:
                continue
            
            content = result.result.data_array[0][0]
            if not content or not content.strip():
                continue
            
            # Extract entities for redaction using NER
            entities_to_redact = extract_entities_for_redaction(content)
            
            if not entities_to_redact:
                redacted_files.append({
                    "original_file": original_file_path,
                    "redacted_file": original_file_path,
                    "entities_count": 0,
                    "status": "no_entities_found",
                    "format": file_format
                })
                continue
            
            # Perform format-specific redaction
            redacted_path = handler.redact_content(original_file_path, entities_to_redact)
            
            redacted_files.append({
                "original_file": original_file_path,
                "redacted_file": redacted_path,
                "entities_count": len(entities_to_redact),
                "entities": entities_to_redact,
                "status": "redacted",
                "format": file_format
            })
        
        return {
            "success": True,
            "message": f"Successfully processed {len(redacted_files)} file(s) for redaction",
            "redacted_files": redacted_files
        }
        
    except Exception as e:
        print(f"Document redaction error: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to redact documents: {str(e)}")

@app.get("/api/download-exported/{filename}")
def download_exported_file(filename: str):
    """Download exported file by filename"""
    if not w:
        raise HTTPException(status_code=500, detail="Databricks connection is not configured.")
    
    if not filename:
        raise HTTPException(status_code=400, detail="filename parameter is required")
    
    # Construct full file path
    file_path = f"{get_uc_volume_path().rstrip('/')}/{filename}"
    
    try:
        print(f"Downloading exported file: {file_path}")
        
        # Download the file from UC Volume
        file_response = w.files.download(file_path=file_path)
        
        # Determine file extension and MIME type
        file_ext = get_file_extension(filename)
        mime_type = get_mime_type(file_ext)
        
        # Create a temporary file to store the downloaded content
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext) as temp_file:
            temp_file_path = temp_file.name
            
            # Handle different response types from Databricks SDK
            response_type = type(file_response).__name__
            print(f"Download response type: {response_type}")
            
            try:
                if hasattr(file_response, 'iter_content') and callable(getattr(file_response, 'iter_content')):
                    for chunk in file_response.iter_content(chunk_size=8192):
                        if chunk:
                            temp_file.write(chunk)
                elif hasattr(file_response, 'content'):
                    content_data = getattr(file_response, 'content')
                    if isinstance(content_data, bytes):
                        temp_file.write(content_data)
                    else:
                        raise Exception(f"Content is not bytes: {type(content_data)}")
                elif hasattr(file_response, 'read') and callable(getattr(file_response, 'read')):
                    content_data = file_response.read()
                    if isinstance(content_data, bytes):
                        temp_file.write(content_data)
                    else:
                        temp_file.write(content_data.encode() if isinstance(content_data, str) else content_data)
                elif isinstance(file_response, bytes):
                    temp_file.write(file_response)
                else:
                    raise Exception(f"Unsupported download response type: {response_type}")
                    
            except Exception as e:
                print(f"Error writing to temp file: {e}")
                os.unlink(temp_file_path)
                raise Exception(f"Failed to process download response: {str(e)}")
        
        # Check if temp file has content
        file_size = os.path.getsize(temp_file_path)
        print(f"Temp file size: {file_size} bytes")
        
        if file_size == 0:
            os.unlink(temp_file_path)
            raise HTTPException(status_code=404, detail="Exported file not found or empty")
        
        # Return the file as a download using FileResponse
        def cleanup_temp_file():
            try:
                os.unlink(temp_file_path)
                print(f"Cleaned up temp file: {temp_file_path}")
            except:
                pass
        
        response = FileResponse(
            path=temp_file_path,
            media_type=mime_type,
            filename=filename
        )
        
        atexit.register(cleanup_temp_file)
        return response
        
    except Exception as e:
        print(f"Exported file download error: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to download exported file: {str(e)}")

# Mount static files for Next.js assets (_next directory, favicon, etc.)
import os
# Use absolute path in Databricks Apps environment
if os.path.exists("/Workspace/Users/q.yu@databricks.com/databricks_apps/pdf-redaction-app/static"):
    target_dir = "/Workspace/Users/q.yu@databricks.com/databricks_apps/pdf-redaction-app/static"
elif os.path.exists("/Workspace/Users/q.yu@databricks.com/databricks_apps/pdf-redaction-app/static"):
    target_dir = "/Workspace/Users/q.yu@databricks.com/databricks_apps/pdf-redaction-app/static"
else:
    # Fallback for local development
    target_dir = "static"

print(f"📁 Serving static files from: {target_dir}")
print(f"📁 _next directory exists: {os.path.exists(f'{target_dir}/_next')}")

# Mount Next.js static assets with proper error handling
try:
    if os.path.exists(f"{target_dir}/_next"):
        app.mount("/_next", StaticFiles(directory=f"{target_dir}/_next"), name="nextjs-assets")
        print("✅ Successfully mounted /_next static files")
    else:
        print("❌ _next directory not found - static assets will not be served")
except Exception as e:
    print(f"❌ Failed to mount static files: {e}")

# Serve other static files with better error handling
@app.get("/favicon.ico")
def favicon():
    try:
        favicon_path = f"{target_dir}/favicon.ico"
        if os.path.exists(favicon_path):
            return FileResponse(favicon_path)
        else:
            print(f"❌ Favicon not found at {favicon_path}")
            raise HTTPException(status_code=404, detail="Favicon not found")
    except Exception as e:
        print(f"❌ Error serving favicon: {e}")
        raise HTTPException(status_code=500, detail="Error serving favicon")

@app.get("/file.svg")  
def file_svg():
    try:
        file_path = f"{target_dir}/file.svg"
        if os.path.exists(file_path):
            return FileResponse(file_path)
        else:
            print(f"❌ file.svg not found at {file_path}")
            raise HTTPException(status_code=404, detail="file.svg not found")
    except Exception as e:
        print(f"❌ Error serving file.svg: {e}")
        raise HTTPException(status_code=500, detail="Error serving file.svg")

# Add a catch-all route for static assets
@app.get("/{asset_path:path}")
def serve_static_asset(asset_path: str):
    """Serve static assets with fallback to main page"""
    # Handle static assets
    if any(asset_path.endswith(ext) for ext in ['.js', '.css', '.woff2', '.svg', '.png', '.ico']):
        static_file_path = f"{target_dir}/{asset_path}"
        if os.path.exists(static_file_path):
            print(f"✅ Serving static asset: {asset_path}")
            return FileResponse(static_file_path)
        else:
            print(f"❌ Static asset not found: {asset_path} at {static_file_path}")
            raise HTTPException(status_code=404, detail=f"Static asset not found: {asset_path}")
    
    # Handle page routes - continue with existing logic
    return serve_react_app(asset_path)

def serve_react_app(full_path: str):
    """Handle Next.js page routes - serve appropriate index.html"""
    # If the request is for a specific HTML file, serve it
    if full_path.endswith('.html'):
        file_path = f"{target_dir}/{full_path}"
        if os.path.exists(file_path):
            return FileResponse(file_path)
    
    
    # For the next-steps route, serve its specific page
    if full_path.startswith("next-steps"):
        file_path = f"{target_dir}/next-steps/index.html"
        if os.path.exists(file_path):
            return FileResponse(file_path)
    
    # For the document-intelligence route, serve its specific page
    if full_path.startswith("document-intelligence"):
        file_path = f"{target_dir}/document-intelligence/index.html"
        if os.path.exists(file_path):
            return FileResponse(file_path)
        
    # For all other routes, serve the main index.html
    return FileResponse(f"{target_dir}/index.html") 